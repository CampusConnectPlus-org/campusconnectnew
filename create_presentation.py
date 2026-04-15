from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

slides = [
    {
        "title": "CampusConnect Project",
        "content": [
            "College campus management platform for alumni, events, placements, and admin oversight.",
            "Built with a React frontend and Node.js + Express backend."
        ]
    },
    {
        "title": "Project Overview",
        "content": [
            "User-facing portal for students, alumni, and administrators.",
            "Admin dashboard supports event management, alumni records, club management, and placement tracking.",
            "Backend APIs provide authentication, data management, and file uploads."
        ]
    },
    {
        "title": "Frontend Overview",
        "content": [
            "React + Vite application.",
            "Uses React Router for navigation across pages: Login, Alumni, Clubs, Events, Placement, Dashboard.",
            "Includes admin pages for managing users, events, placements, and alumni data."
        ]
    },
    {
        "title": "Backend Overview",
        "content": [
            "Node.js and Express server running on port 5000.",
            "MongoDB Atlas database with Mongoose models for Admin, User, Alumni, Event, Placement entities.",
            "API routes: /api/auth, /admin, /api/placements, /api/events, /api/clubs."
        ]
    },
    {
        "title": "Key Features",
        "content": [
            "Authentication and role-based admin access.",
            "Event creation and listings.",
            "Alumni directory and profile management.",
            "Placement drives, statistics, and trends.",
            "Club management and activity tracking."
        ]
    },
    {
        "title": "Core Components",
        "content": [
            "AdminNavbar, Sidebar, UsersTable, RecentEvents, ActivityChart.",
            "AddAlumniModal, AlumniSection, CompanySlider.",
            "Protected routes for secure access control."
        ]
    },
    {
        "title": "Technology Stack",
        "content": [
            "Frontend: React, Vite, React Router, Axios, framer-motion.",
            "Backend: Express, MongoDB Atlas, Mongoose, JWT, bcryptjs, multer.",
            "Development: ESLint, Node.js, JavaScript modules."
        ]
    },
    {
        "title": "Architecture",
        "content": [
            "Frontend client consumes REST APIs from the backend.",
            "Backend serves API requests, interacts with MongoDB, and handles authentication and uploads.",
            "Static assets and uploaded files are served from Express endpoints."
        ]
    },
    {
        "title": "Demo / Next Steps",
        "content": [
            "Show user login and admin dashboard flow.",
            "Demonstrate event creation and alumni listing.",
            "Future enhancements: role-based access, analytics dashboard, mobile responsiveness." 
        ]
    }
]

for slide_data in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = slide_data["title"]
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for index, line in enumerate(slide_data["content"]):
        if index == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 1
        p.font.size = Pt(24)

prs.save("CampusConnect_Project_Presentation.pptx")
print("Created CampusConnect_Project_Presentation.pptx")
